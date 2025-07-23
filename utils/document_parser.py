"""
Complete Document Parser with All File Type Support

This module provides comprehensive document parsing for multiple file formats.
"""

import os
import logging
import re
import csv
from typing import List, Optional, Dict, Any
from dataclasses import dataclass

logger = logging.getLogger(__name__)

# Try to import optional dependencies
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    logger.warning("pandas not available - CSV processing will be limited")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX files not supported")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPTX files not supported")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyPDF2 not available - PDF files not supported")

@dataclass
class ChunkMetadata:
    source_file: str
    chunk_index: int
    chunk_type: str
    page_number: Optional[int] = None
    section_title: Optional[str] = None
    word_count: int = 0
    char_count: int = 0

class DocumentParser:
    # Configuration for chunking
    DEFAULT_CHUNK_SIZE = 1000  # characters
    DEFAULT_CHUNK_OVERLAP = 200  # characters
    MIN_CHUNK_SIZE = 100  # minimum chunk size
    
    @staticmethod
    def parse_file(file_path: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """Parse a file and return chunks with improved chunking strategy"""
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return []
        
        chunk_size = chunk_size or DocumentParser.DEFAULT_CHUNK_SIZE
        chunk_overlap = chunk_overlap or DocumentParser.DEFAULT_CHUNK_OVERLAP
        
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.pdf' and PDF_AVAILABLE:
                raw_chunks = DocumentParser._parse_pdf(file_path)
            elif ext == '.docx' and DOCX_AVAILABLE:
                raw_chunks = DocumentParser._parse_docx(file_path)
            elif ext == '.pptx' and PPTX_AVAILABLE:
                raw_chunks = DocumentParser._parse_pptx(file_path)
            elif ext == '.csv':
                raw_chunks = DocumentParser._parse_csv(file_path)
            elif ext in ['.txt', '.md']:
                raw_chunks = DocumentParser._parse_txt(file_path)
            else:
                # Fallback: try to read as text
                raw_chunks = DocumentParser._parse_as_text(file_path)
            
            # Apply intelligent chunking
            processed_chunks = DocumentParser._apply_intelligent_chunking(
                raw_chunks, chunk_size, chunk_overlap
            )
            
            logger.info(f"Parsed {file_path}: {len(raw_chunks)} raw chunks -> {len(processed_chunks)} processed chunks")
            return processed_chunks
            
        except Exception as e:
            logger.error(f"Error parsing file {file_path}: {str(e)}", exc_info=True)
            # Return a placeholder chunk with error info
            return [f"[Error parsing {os.path.basename(file_path)}: {str(e)}]"]
    
    @staticmethod
    def _apply_intelligent_chunking(raw_chunks: List[str], chunk_size: int, chunk_overlap: int) -> List[str]:
        """Apply intelligent chunking strategy to raw text chunks"""
        if not raw_chunks:
            return []
        
        # Combine all raw chunks into a single text
        full_text = "\n\n".join(raw_chunks)
        
        # Split into sentences for better chunk boundaries
        sentences = DocumentParser._split_into_sentences(full_text)
        
        chunks = []
        current_chunk = ""
        current_size = 0
        
        for sentence in sentences:
            sentence_size = len(sentence)
            
            # If adding this sentence would exceed chunk size
            if current_size + sentence_size > chunk_size and current_chunk:
                # Save current chunk if it meets minimum size
                if current_size >= DocumentParser.MIN_CHUNK_SIZE:
                    chunks.append(current_chunk.strip())
                
                # Start new chunk with overlap
                if chunk_overlap > 0 and current_chunk:
                    overlap_text = DocumentParser._get_overlap_text(current_chunk, chunk_overlap)
                    current_chunk = overlap_text + " " + sentence
                    current_size = len(current_chunk)
                else:
                    current_chunk = sentence
                    current_size = sentence_size
            else:
                # Add sentence to current chunk
                if current_chunk:
                    current_chunk += " " + sentence
                else:
                    current_chunk = sentence
                current_size = len(current_chunk)
        
        # Add final chunk
        if current_chunk.strip() and len(current_chunk) >= DocumentParser.MIN_CHUNK_SIZE:
            chunks.append(current_chunk.strip())
        
        return chunks if chunks else [full_text[:chunk_size]]  # Fallback
    
    @staticmethod
    def _split_into_sentences(text: str) -> List[str]:
        """Split text into sentences using regex"""
        sentence_pattern = r'(?<=[.!?])\s+'
        sentences = re.split(sentence_pattern, text)
        return [s.strip() for s in sentences if s.strip()]
    
    @staticmethod
    def _get_overlap_text(text: str, overlap_size: int) -> str:
        """Get the last part of text for overlap"""
        if len(text) <= overlap_size:
            return text
        return text[-overlap_size:]
    
    @staticmethod
    def _parse_pdf(file_path: str) -> List[str]:
        """Parse PDF files"""
        chunks = []
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        text = DocumentParser._clean_text(text)
                        chunks.append(f"[Page {page_num}]\n{text}")
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {str(e)}")
            return [f"[PDF parsing error: {str(e)}]"]
        
        return chunks
    
    @staticmethod
    def _parse_docx(file_path: str) -> List[str]:
        """Parse DOCX files"""
        chunks = []
        try:
            doc = Document(file_path)
            current_section = ""
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # Check if this looks like a heading
                if DocumentParser._is_heading(para):
                    if current_section:
                        chunks.append(DocumentParser._clean_text(current_section))
                    current_section = f"[Heading] {text}\n"
                else:
                    current_section += text + "\n"
            
            # Add final section
            if current_section:
                chunks.append(DocumentParser._clean_text(current_section))
                
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {str(e)}")
            return [f"[DOCX parsing error: {str(e)}]"]
        
        return chunks
    
    @staticmethod
    def _is_heading(paragraph) -> bool:
        """Check if a paragraph is likely a heading"""
        if hasattr(paragraph, 'style') and paragraph.style:
            style_name = paragraph.style.name.lower()
            return 'heading' in style_name or 'title' in style_name
        return False
    
    @staticmethod
    def _parse_pptx(file_path: str) -> List[str]:
        """Parse PowerPoint files"""
        chunks = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"[Slide {slide_num}]\n"
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content += "\n".join(slide_text)
                    chunks.append(DocumentParser._clean_text(slide_content))
                    
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {str(e)}")
            return [f"[PPTX parsing error: {str(e)}]"]
        
        return chunks
    
    @staticmethod
    def _parse_csv(file_path: str) -> List[str]:
        """Parse CSV files with or without pandas"""
        chunks = []
        
        if PANDAS_AVAILABLE:
            try:
                df = pd.read_csv(file_path)
                
                # Add metadata about the CSV
                chunks.append(f"[CSV Metadata]\nColumns: {', '.join(df.columns)}\nRows: {len(df)}")
                
                # Process each column
                for col in df.columns:
                    col_info = f"[Column: {col}]\n"
                    col_info += f"Data type: {df[col].dtype}\n"
                    col_info += f"Non-null values: {df[col].count()}\n"
                    
                    # Add sample values
                    sample_values = df[col].dropna().head(10).tolist()
                    col_info += f"Sample values: {', '.join(map(str, sample_values))}"
                    
                    chunks.append(col_info)
                
                # Add summary statistics for numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    stats_info = "[Numeric Statistics]\n"
                    stats_info += df[numeric_cols].describe().to_string()
                    chunks.append(stats_info)
                    
            except Exception as e:
                logger.error(f"Error parsing CSV with pandas {file_path}: {str(e)}")
                # Fallback to basic CSV parsing
                chunks = DocumentParser._parse_csv_basic(file_path)
        else:
            chunks = DocumentParser._parse_csv_basic(file_path)
        
        return chunks
    
    @staticmethod
    def _parse_csv_basic(file_path: str) -> List[str]:
        """Basic CSV parsing without pandas"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                # Try to detect delimiter
                sample = f.read(1024)
                f.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.reader(f, delimiter=delimiter)
                rows = list(reader)
                
                if not rows:
                    return ["[Empty CSV file]"]
                
                # Add metadata about the CSV
                headers = rows[0] if rows else []
                data_rows = rows[1:] if len(rows) > 1 else []
                
                chunks.append(f"[CSV Metadata]\nColumns: {', '.join(headers)}\nRows: {len(data_rows)}")
                
                # Process headers
                if headers:
                    chunks.append(f"[CSV Headers]\n{', '.join(headers)}")
                
                # Process data in batches
                batch_size = 10
                for i in range(0, len(data_rows), batch_size):
                    batch = data_rows[i:i+batch_size]
                    batch_text = f"[CSV Data Batch {i//batch_size + 1}]\n"
                    
                    for row in batch:
                        if len(row) == len(headers):
                            row_text = []
                            for header, value in zip(headers, row):
                                if value.strip():
                                    row_text.append(f"{header}: {value}")
                            if row_text:
                                batch_text += " | ".join(row_text) + "\n"
                    
                    if batch_text.strip():
                        chunks.append(batch_text)
                
        except Exception as e:
            logger.error(f"Error parsing CSV {file_path}: {str(e)}")
            return [f"[CSV parsing error: {str(e)}]"]
        
        return chunks
    
    @staticmethod
    def _parse_txt(file_path: str) -> List[str]:
        """Parse text files with encoding detection"""
        chunks = []
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                # Split by double newlines (paragraphs) or single newlines for smaller files
                if len(content) > 5000:
                    chunks = [chunk.strip() for chunk in content.split('\n\n') if chunk.strip()]
                else:
                    chunks = [chunk.strip() for chunk in content.split('\n') if chunk.strip()]
                
                # If we got here, encoding worked
                break
                
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error parsing text file {file_path}: {str(e)}")
                break
        
        return chunks if chunks else [f"[Could not read text file: {file_path}]"]
    
    @staticmethod
    def _parse_as_text(file_path: str) -> List[str]:
        """Fallback: try to parse any file as text"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return [content] if content.strip() else [f"[Empty or unreadable file: {os.path.basename(file_path)}]"]
        except Exception as e:
            logger.error(f"Error reading file as text {file_path}: {str(e)}")
            return [f"[Unreadable file: {os.path.basename(file_path)}]"]
    
    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean and normalize text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters that might cause issues
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', '', text)
        
        return text.strip()
    
    @staticmethod
    def get_supported_extensions() -> List[str]:
        """Get list of supported file extensions"""
        extensions = ['.txt', '.md']
        
        if PDF_AVAILABLE:
            extensions.append('.pdf')
        if DOCX_AVAILABLE:
            extensions.append('.docx')
        if PPTX_AVAILABLE:
            extensions.append('.pptx')
        
        extensions.append('.csv')  # Always supported (basic or pandas)
        
        return extensions
    
    @staticmethod
    def is_supported_file(file_path: str) -> bool:
        """Check if file type is supported"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in DocumentParser.get_supported_extensions() or os.path.exists(file_path)